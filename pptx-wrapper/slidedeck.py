#!/usr/local/bin/python3
from pptx import Presentation
from pptx.util import Inches

# height = Inches(0.8)
# img = sld.shapes.add_picture(fn[0], 0, 0)
# img_path = fn[0]
# Inches(5)
# Inches(prs.slide_height)
# left = top = Inches(1)
# left = top = Inches(2.0)

# pic = slide.shapes.add_picture(img_path, Inches(5), top, height=Inches(5.5))


# prs = Presentation('template.pptx')

# title_placeholder = slide.shapes.title
# title_placeholder.text = 'Air-speed Velocity of Unladen Swallows'


def available_layouts(prs):
    layouts = [prs.slide_layouts[n].name for n in range(
        len(prs.slide_layouts))]
    return layouts


def new_slide(prs, layout='Blank'):
    layouts = available_layouts(prs)
    if layout not in layouts:
        print('No layout by that name! Available layouts:\n')
        print(layouts)
        return
    slide = prs.slides.add_slide(prs.slide_layouts[layouts.index(layout)])
    for shape in slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))
    return slide


# def save(prs, outname=None):

#     # shape shapes – auto shapes with fill and an outline
#     # text boxes – auto shapes with no fill and no outline
#     # placeholders – auto shapes that can appear on a slide layout or master and be inherited on slides that use that layout, allowing content to be added that takes on the formatting of the placeholder
#     # picture – as described above
#     # table – that row and column thing

#     prs.save('test.pptx')


def df_to_table(slide, df, pos=[1.5, 0.25, 9.25, 6.0], colnames=None):
    """Converts a Pandas DataFrame to a PowerPoint table on the given
    Slide of a PowerPoint presentation.
    The table is a standard Powerpoint table, and can easily be modified with the Powerpoint tools,
    for example: resizing columns, changing formatting etc.
    Arguments:
     - slide: slide object from the python-pptx library containing the slide on which you want the table to appear
     - df: Pandas DataFrame with the data
    Optional arguments:
     - colnames
     https://github.com/robintw/PandasToPowerpoint/blob/master/PandasToPowerpoint.py
     """
    rows, cols = df.shape
    top = Inches(pos[0])
    left = Inches(pos[1])
    width = Inches(pos[2])
    height = Inches(pos[3])
    res = slide.shapes.add_table(rows + 1, cols, left, top, width, height)

    if colnames is None:
        colnames = list(df.columns)

    # Insert the column names
    for col_index, col_name in enumerate(colnames):
        # Column names can be tuples
        if not isinstance(col_name, str):
            col_name = " ".join(col_name)
        res.table.cell(0, col_index).text = col_name

    m = df.values()

    for row in range(rows):
        for col in range(cols):
            val = m[row, col]
            text = str(val)
            res.table.cell(row + 1, col).text = text


def create_chart(df, filename):
    """ Create a simple bar chart saved to the filename based on the dataframe
    passed to the function
    """
    df['total'] = df['Quantity'] * df['Price']
    final_plot = df.groupby('Name')['total'].sum().order().plot(kind='barh')
    fig = final_plot.get_figure()
    fig.set_size_inches(6, 4.5)
    fig.savefig(filename, bbox_inches='tight', dpi=600)


def create_ppt(input, output, report_data, chart):
    """ Take the input powerpoint file and use it as the template for the output
    file.
    """
    prs = Presentation(input)
    # Use the output from analyze_ppt to understand which layouts and placeholders
    # to use
    # Create a title slide first
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Quarterly Report"
    subtitle.text = "Generated on {:%m-%d-%Y}".format(date.today())

    prs.save(output)
